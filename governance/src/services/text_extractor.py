"""テキスト抽出 — ファイルコンテンツからプレーンテキストを抽出する

詳細設計 4.4 節準拠

対応フォーマット: docx, xlsx, pptx, pdf, csv, txt, markdown
"""

from __future__ import annotations

import csv
import io
from typing import Callable

from shared.config import SSM_MAX_TEXT_LENGTH, get_ssm_int
from shared.logger import get_logger

logger = get_logger(__name__)


def extract_text(file_content: bytes, mime_type: str) -> str:
    """ファイルコンテンツからテキストを抽出するディスパッチャ。

    対応形式でない場合や抽出失敗時は空文字列を返す。
    """
    extractor = TEXT_EXTRACTORS.get(mime_type)
    if extractor is None:
        logger.info(f"Unsupported mime_type for text extraction: {mime_type}")
        return ""

    try:
        return extractor(file_content)
    except Exception as e:
        logger.warning(f"Text extraction failed for {mime_type}: {e}")
        return ""


def extract_docx(content: bytes) -> str:
    """docx からテキストを抽出する（段落 + テーブル）。"""
    from docx import Document

    doc = Document(io.BytesIO(content))
    texts: list[str] = []
    for paragraph in doc.paragraphs:
        if paragraph.text:
            texts.append(paragraph.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text:
                    texts.append(cell.text)
    return "\n".join(texts)


def extract_xlsx(content: bytes) -> str:
    """xlsx からテキストを抽出する（全シート・全セル）。"""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    texts: list[str] = []
    try:
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                for cell_value in row:
                    if cell_value is not None:
                        texts.append(str(cell_value))
    finally:
        wb.close()
    return "\n".join(texts)


def extract_pptx(content: bytes) -> str:
    """pptx からテキストを抽出する（スライド + ノート）。"""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text:
                        texts.append(paragraph.text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
                if paragraph.text:
                    texts.append(paragraph.text)
    return "\n".join(texts)


def extract_pdf(content: bytes) -> str:
    """pdf からテキストを抽出する（ページごと）。"""
    from PyPDF2 import PdfReader

    reader = PdfReader(io.BytesIO(content))
    texts: list[str] = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            texts.append(page_text)
    return "\n".join(texts)


def extract_csv_text(content: bytes) -> str:
    """csv からテキストを抽出する（全セル）。"""
    text = content.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    texts: list[str] = []
    for row in reader:
        for cell in row:
            if cell:
                texts.append(cell)
    return "\n".join(texts)


def extract_plain(content: bytes) -> str:
    """プレーンテキストとしてそのまま読み込む。"""
    return content.decode("utf-8", errors="replace")


def truncate_text(text: str, max_length: int | None = None) -> str:
    """テキストを指定文字数に制限する。

    max_length が None の場合は SSM パラメータから取得する。
    """
    if max_length is None:
        max_length = _get_max_text_length()
    if len(text) <= max_length:
        return text
    return text[:max_length]


# MIME Type → 抽出関数のレジストリ
TEXT_EXTRACTORS: dict[str, Callable[[bytes], str]] = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": extract_docx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": extract_xlsx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": extract_pptx,
    "application/pdf": extract_pdf,
    "text/csv": extract_csv_text,
    "text/plain": extract_plain,
    "text/markdown": extract_plain,
}

# サポートされる MIME Type の集合（外部から参照可能）
SUPPORTED_MIME_TYPES = frozenset(TEXT_EXTRACTORS.keys())


def is_supported_format(mime_type: str) -> bool:
    """テキスト抽出に対応した MIME Type かどうかを判定する。"""
    return mime_type in SUPPORTED_MIME_TYPES


def _get_max_text_length() -> int:
    try:
        return get_ssm_int(SSM_MAX_TEXT_LENGTH, default=500000)
    except Exception:
        return 500000
