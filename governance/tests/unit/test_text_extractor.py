"""text_extractor 単体テスト

テキスト抽出の各フォーマット対応・エッジケースを検証する。
サンプルファイルの生成にはそれぞれのライブラリを使用する。
"""

import io

import pytest


class TestExtractPlain:
    def test_utf8_text(self):
        from services.text_extractor import extract_plain

        content = "Hello, World!\n日本語テスト".encode("utf-8")
        result = extract_plain(content)
        assert "Hello, World!" in result
        assert "日本語テスト" in result

    def test_empty_bytes(self):
        from services.text_extractor import extract_plain

        result = extract_plain(b"")
        assert result == ""


class TestExtractCsv:
    def test_basic_csv(self):
        from services.text_extractor import extract_csv_text

        content = "名前,電話番号\n田中太郎,03-1234-5678\n".encode("utf-8")
        result = extract_csv_text(content)
        assert "田中太郎" in result
        assert "03-1234-5678" in result

    def test_empty_csv(self):
        from services.text_extractor import extract_csv_text

        result = extract_csv_text(b"")
        assert result == ""


class TestExtractDocx:
    def test_basic_docx(self):
        from docx import Document as DocxDocument

        from services.text_extractor import extract_docx

        doc = DocxDocument()
        doc.add_paragraph("テスト文書です。")
        doc.add_paragraph("個人情報を含みます。")
        buf = io.BytesIO()
        doc.save(buf)

        result = extract_docx(buf.getvalue())
        assert "テスト文書です。" in result
        assert "個人情報を含みます。" in result

    def test_docx_with_table(self):
        from docx import Document as DocxDocument

        from services.text_extractor import extract_docx

        doc = DocxDocument()
        doc.add_paragraph("概要")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "名前"
        table.cell(0, 1).text = "電話"
        table.cell(1, 0).text = "田中"
        table.cell(1, 1).text = "090-1234"
        buf = io.BytesIO()
        doc.save(buf)

        result = extract_docx(buf.getvalue())
        assert "概要" in result
        assert "田中" in result
        assert "090-1234" in result

    def test_empty_docx(self):
        from docx import Document as DocxDocument

        from services.text_extractor import extract_docx

        doc = DocxDocument()
        buf = io.BytesIO()
        doc.save(buf)

        result = extract_docx(buf.getvalue())
        assert result == ""


class TestExtractXlsx:
    def test_basic_xlsx(self):
        from openpyxl import Workbook

        from services.text_extractor import extract_xlsx

        wb = Workbook()
        ws = wb.active
        ws["A1"] = "名前"
        ws["B1"] = "給与"
        ws["A2"] = "田中太郎"
        ws["B2"] = 500000
        buf = io.BytesIO()
        wb.save(buf)

        result = extract_xlsx(buf.getvalue())
        assert "田中太郎" in result
        assert "500000" in result

    def test_multi_sheet_xlsx(self):
        from openpyxl import Workbook

        from services.text_extractor import extract_xlsx

        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Sheet1"
        ws1["A1"] = "シート1データ"
        ws2 = wb.create_sheet("Sheet2")
        ws2["A1"] = "シート2データ"
        buf = io.BytesIO()
        wb.save(buf)

        result = extract_xlsx(buf.getvalue())
        assert "シート1データ" in result
        assert "シート2データ" in result

    def test_empty_xlsx(self):
        from openpyxl import Workbook

        from services.text_extractor import extract_xlsx

        wb = Workbook()
        buf = io.BytesIO()
        wb.save(buf)

        result = extract_xlsx(buf.getvalue())
        assert result == ""


class TestExtractPptx:
    def test_basic_pptx(self):
        from pptx import Presentation

        from services.text_extractor import extract_pptx

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "プレゼンテーション"
        slide.placeholders[1].text = "機密情報含む"
        buf = io.BytesIO()
        prs.save(buf)

        result = extract_pptx(buf.getvalue())
        assert "プレゼンテーション" in result
        assert "機密情報含む" in result

    def test_empty_pptx(self):
        from pptx import Presentation

        from services.text_extractor import extract_pptx

        prs = Presentation()
        buf = io.BytesIO()
        prs.save(buf)

        result = extract_pptx(buf.getvalue())
        assert result == ""


class TestExtractPdf:
    def test_basic_pdf(self):
        from PyPDF2 import PdfWriter

        from services.text_extractor import extract_pdf

        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        buf = io.BytesIO()
        writer.write(buf)

        result = extract_pdf(buf.getvalue())
        assert isinstance(result, str)


class TestExtractText:
    def test_dispatch_plain(self):
        from services.text_extractor import extract_text

        result = extract_text(b"hello world", "text/plain")
        assert result == "hello world"

    def test_dispatch_markdown(self):
        from services.text_extractor import extract_text

        result = extract_text(b"# heading\ncontent", "text/markdown")
        assert "heading" in result

    def test_unsupported_format(self):
        from services.text_extractor import extract_text

        result = extract_text(b"\x00\x01\x02", "application/octet-stream")
        assert result == ""

    def test_invalid_content_returns_empty(self):
        from services.text_extractor import extract_text

        result = extract_text(b"not a docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        assert result == ""


class TestTruncateText:
    def test_within_limit(self):
        from services.text_extractor import truncate_text

        result = truncate_text("short text", max_length=100)
        assert result == "short text"

    def test_exceeds_limit(self):
        from services.text_extractor import truncate_text

        text = "a" * 1000
        result = truncate_text(text, max_length=500)
        assert len(result) == 500

    def test_exact_limit(self):
        from services.text_extractor import truncate_text

        text = "a" * 100
        result = truncate_text(text, max_length=100)
        assert len(result) == 100


class TestIsSupportedFormat:
    def test_supported_formats(self):
        from services.text_extractor import is_supported_format

        assert is_supported_format("text/plain") is True
        assert is_supported_format("text/csv") is True
        assert is_supported_format("text/markdown") is True
        assert is_supported_format("application/pdf") is True
        assert is_supported_format(
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ) is True
        assert is_supported_format(
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ) is True
        assert is_supported_format(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ) is True

    def test_unsupported_formats(self):
        from services.text_extractor import is_supported_format

        assert is_supported_format("application/octet-stream") is False
        assert is_supported_format("image/png") is False
        assert is_supported_format("") is False
